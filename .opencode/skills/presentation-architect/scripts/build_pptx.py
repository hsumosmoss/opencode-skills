#!/usr/bin/env python3
import sys
import yaml
import os
import glob
import re
import tempfile
import zipfile
import shutil
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image

def hex_to_rgb(h):
    h = h.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def calculate_fit_dimensions(img_path, max_width_in, max_height_in):
    try:
        with Image.open(img_path) as img:
            w_px, h_px = img.size
        aspect_ratio = w_px / float(h_px)
        box_aspect = max_width_in / float(max_height_in)
        if aspect_ratio > box_aspect:
            return max_width_in, max_width_in / aspect_ratio
        else:
            return max_height_in * aspect_ratio, max_height_in
    except Exception as e:
        print(f"Error reading image {img_path}: {e}")
        return max_width_in, max_height_in

def set_font_style(paragraph, font_name="Microsoft JhengHei", size=24, color_hex="#333333", bold=False):
    paragraph.font.name = font_name
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*hex_to_rgb(color_hex))
    paragraph.space_after = Pt(8)
    paragraph.line_spacing = 1.3

def place_image(slide, img_path, box_x, box_y, box_w, box_h):
    fit_w, fit_h = calculate_fit_dimensions(img_path, box_w, box_h)
    center_x = box_x + (box_w - fit_w) / 2
    center_y = box_y + (box_h - fit_h) / 2
    slide.shapes.add_picture(img_path, Inches(center_x), Inches(center_y), width=Inches(fit_w), height=Inches(fit_h))

def place_multiple_images(slide, img_paths, box_x, box_y, box_w, box_h):
    count = len(img_paths)
    if count == 0:
        return
    elif count == 1:
        place_image(slide, img_paths[0], box_x, box_y, box_w, box_h)
    elif count == 2:
        half_h = box_h / 2.0
        place_image(slide, img_paths[0], box_x, box_y, box_w, half_h)
        place_image(slide, img_paths[1], box_x, box_y + half_h, box_w, half_h)
    else:
        half_w = box_w / 2.0
        half_h = box_h / 2.0
        coords = [(0,0), (1,0), (0,1), (1,1)]
        for i, img in enumerate(img_paths[:4]):
            col, row = coords[i]
            place_image(slide, img, box_x + col*half_w, box_y + row*half_h, half_w, half_h)

def extract_docx_images(docx_path, out_dir):
    """
    從 DOCX 深度解析：將圖片對應到 '第X頁' 的標題，並解壓縮出對應的圖片檔。
    """
    page_images_map = {}
    
    with zipfile.ZipFile(docx_path, 'r') as docx:
        # 解析 rels 映射
        rels_xml = docx.read('word/_rels/document.xml.rels')
        soup_rels = BeautifulSoup(rels_xml, 'xml')
        rid_to_target = {}
        for rel in soup_rels.find_all('Relationship'):
            rid_to_target[rel.get('Id')] = rel.get('Target')
            
        # 解析 document.xml
        doc_xml = docx.read('word/document.xml')
        soup_doc = BeautifulSoup(doc_xml, 'xml')
        
        current_page = 0
        
        for p in soup_doc.find_all('w:p'):
            text = "".join(t.text for t in p.find_all('w:t'))
            if text:
                match = re.search(r'第\s*(\d+)\s*頁', text)
                if match:
                    current_page = int(match.group(1))
                    if current_page not in page_images_map:
                        page_images_map[current_page] = []
                        
            # 尋找圖片
            for blip in p.find_all('a:blip'):
                embed_id = blip.get('r:embed')
                if embed_id and embed_id in rid_to_target:
                    target_path = rid_to_target[embed_id]
                    if target_path.startswith('media/'):
                        internal_path = f"word/{target_path}"
                    else:
                        internal_path = f"word/{target_path}"
                        
                    ext = os.path.splitext(target_path)[1]
                    out_filename = f"page_{current_page}_{embed_id}{ext}"
                    out_filepath = os.path.join(out_dir, out_filename)
                    
                    try:
                        with open(out_filepath, 'wb') as f:
                            f.write(docx.read(internal_path))
                        
                        if current_page not in page_images_map:
                            page_images_map[current_page] = []
                        page_images_map[current_page].append(out_filepath)
                    except Exception as e:
                        print(f"Failed to extract {internal_path}: {e}")
                        
    return page_images_map

def main():
    if len(sys.argv) < 4:
        print("Usage: python3 build_pptx.py <yaml_path> <image_dir_path> <docx_path>")
        sys.exit(1)
        
    yaml_path = sys.argv[1]
    img_dir = sys.argv[2]
    docx_path = sys.argv[3]
    
    with open(yaml_path, 'r', encoding='utf-8') as f:
        data = yaml.safe_load(f)
        
    config = data.get('presentation_anchored_v5_section_split', {})
    metadata = config.get('metadata', {})
    style = metadata.get('style', {})
    
    primary_hex = style.get('color_primary', "#4285F4")
    bg_hex = style.get('color_background', "#FFFFFF")
    text_hex = style.get('color_text', "#333333")
    
    slides = config.get('slide_execution', [])
    
    # 💡 啟動原生 DOCX 解析引擎
    print("🔍 啟動 DOCX 原生解析引擎...")
    temp_extract_dir = tempfile.mkdtemp(prefix="docx_images_")
    page_to_images = extract_docx_images(docx_path, temp_extract_dir)
    print(f"✅ DOCX 解析完成，找到原圖的頁面有: {list(page_to_images.keys())}")
    
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    blank_slide_layout = prs.slide_layouts[6]
    
    for page_data in slides:
        slide_idx = page_data.get('page', 1)
        elements = page_data.get('elements_manifest', [])
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*hex_to_rgb(bg_hex))
        
        title_text = ""
        body_text = ""
        wants_src_image = False
        
        for el in elements:
            if el['type'] == 'title':
                title_text = el['content']
            elif el['type'] == 'text_block':
                body_text += el['content'] + "\n"
            elif el['type'] == 'image_anchor':
                wants_src_image = True
                
        # 💡 精準掛載該頁的原圖
        src_imgs = []
        if wants_src_image and slide_idx in page_to_images:
            src_imgs = page_to_images[slide_idx]
            
        ip_pattern = os.path.join(img_dir, f"slide_{slide_idx:02d}_ip_*.png")
        ip_imgs = glob.glob(ip_pattern)
        if not ip_imgs:
            ip_pattern2 = os.path.join(img_dir, f"slide_{slide_idx}_ip_*.png")
            ip_imgs = glob.glob(ip_pattern2)
            
        has_src = len(src_imgs) > 0
        has_ip = len(ip_imgs) > 0
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(15), Inches(1.2))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        set_font_style(p, size=40, color_hex=primary_hex, bold=True)
        
        text_len = len(body_text)
        
        if has_src and has_ip:
            if text_len > 120:
                text_box = (0.5, 1.8, 15.0, 2.5)
                src_box = (2.0, 4.5, 5.0, 4.0)
                ip_box = (9.0, 4.5, 5.0, 4.0)
                font_size = 18
            else:
                text_box = (0.5, 2.0, 4.5, 6.5)
                src_box = (5.2, 2.0, 6.5, 6.5)
                ip_box = (12.0, 2.0, 3.5, 6.5)
                font_size = 22
        elif not has_src and has_ip:
            text_box = (0.5, 2.0, 10.0, 6.5)
            src_box = None
            ip_box = (11.0, 2.0, 4.5, 6.5)
            font_size = 22 if text_len <= 200 else 18
        elif has_src and not has_ip:
            text_box = (0.5, 2.0, 6.0, 6.5)
            src_box = (7.0, 2.0, 8.5, 6.5)
            ip_box = None
            font_size = 22 if text_len <= 150 else 18
        else:
            text_box = (0.5, 2.0, 15.0, 6.5)
            src_box = None
            ip_box = None
            font_size = 22 if text_len <= 300 else 18
            
        txBox2 = slide.shapes.add_textbox(Inches(text_box[0]), Inches(text_box[1]), Inches(text_box[2]), Inches(text_box[3]))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = body_text.strip()
        set_font_style(p2, size=font_size, color_hex=text_hex)
        
        if has_src and src_box:
            place_multiple_images(slide, src_imgs, *src_box)
            
        if has_ip and ip_box:
            place_image(slide, ip_imgs[0], *ip_box)

    out_file = os.path.join(img_dir, "Anchored_Presentation.pptx")
    prs.save(out_file)
    print(f"✅ Dynamic Styled PPTX successfully built at {out_file}")
    
    shutil.rmtree(temp_extract_dir)

if __name__ == "__main__":
    main()
