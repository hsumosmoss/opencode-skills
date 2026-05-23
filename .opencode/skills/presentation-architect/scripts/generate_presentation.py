#!/usr/bin/env python3
import sys
import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_presentation(img_dir, output_path, slides_data, bg_hex="0A0C10", text_color_hex="FFFFFF"):
    """
    slides_data is a list of dicts:
    [{"img_pattern": "slide_01*.png", "title": "...", "subtitle": "...", "desc": "..."}]
    """
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    blank_slide_layout = prs.slide_layouts[6]
    
    # Convert hex to RGB
    def hex_to_rgb(h):
        h = h.lstrip('#')
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))
    
    bg_r, bg_g, bg_b = hex_to_rgb(bg_hex)
    t_r, t_g, t_b = hex_to_rgb(text_color_hex)

    for i, data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # 設置背景色
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(bg_r, bg_g, bg_b)
        
        # 交錯排版：偶數頁圖片在左，奇數頁圖片在右
        img_left = Inches(0) if i % 2 == 0 else Inches(7)
        text_left = Inches(9.5) if i % 2 == 0 else Inches(1)
        
        # 尋找圖片
        search_pattern = os.path.join(img_dir, data.get("img_pattern", "*.png"))
        img_paths = glob.glob(search_pattern)
        
        if img_paths:
            # 加入圖片，高度固定 9 吋（如果是 1:1 圖片，寬度也會是 9 吋）
            slide.shapes.add_picture(img_paths[0], img_left, Inches(0), height=Inches(9))
        
        # 加入文字介紹
        txBox = slide.shapes.add_textbox(text_left, Inches(2.5), Inches(5.5), Inches(4))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = data.get("title", "")
        p.font.bold = True
        p.font.size = Pt(50)
        p.font.color.rgb = RGBColor(t_r, t_g, t_b)
        
        if data.get("subtitle"):
            p2 = tf.add_paragraph()
            p2.text = data["subtitle"]
            p2.font.bold = False
            p2.font.size = Pt(32)
            p2.font.color.rgb = RGBColor(0, 255, 255) # 霓虹藍強調色
            
        if data.get("desc"):
            p_space = tf.add_paragraph()
            p_space.text = ""
            p_space.font.size = Pt(14)
            
            p3 = tf.add_paragraph()
            p3.text = data["desc"]
            p3.font.bold = False
            p3.font.size = Pt(24)
            p3.font.color.rgb = RGBColor(180, 190, 200)

    prs.save(output_path)
    print(f"✅ PPTX generated successfully at {output_path}")

# Example usage when run directly (for testing)
if __name__ == "__main__":
    print("This script is meant to be imported and used by the agent, or customized before running.")
